#!/usr/bin/python python3

import requests
import pandas as pd
from datetime import *
from bs4 import BeautifulSoup
from pptx import Presentation
from datetime import datetime
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx import Presentation
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

page = requests.get('https://www.astro.columbia.edu/events/colloquia')
contents = page.content
soup = BeautifulSoup(page.content, 'html.parser')

event = soup.find_all("a", href=lambda href: href and "event?" in href)
f = open('./astrokiosk.pptx',"rb")
prs = Presentation(f)

for i in range(0,len(event)):
    link = "http://www.astro.columbia.edu" + event[i].get('href')
    #print(link)
    page = requests.get(link)
    contents = page.content
    event_soup = BeautifulSoup(page.content, 'html.parser')
    p = event_soup.find_all("p")
    abstract = p[0].get_text()
    #print("abstract:",abstract)
    #print(event_soup.prettify)
    #long_title = event[i].get('title')
    specs = event_soup.find("ul",class_="specs").get_text().strip()
    specs = specs.replace('Date:','')
    specs = specs.replace('Location:','Location: ')
    specs = specs.replace('Speaker:','')
    specs = specs.replace('Time:','')
    specs = specs.replace('Host:','Host: ')
    speaker = event[i].find(class_="speaker").get_text().strip()
    speaker = speaker.replace('Speaker:','')
    #title = event[i].find(class_="title").get_text().strip()
    title = event_soup.find_all("div",class_="content event")
    title = title[-1].find_all("h3")[0].get_text()
    date = event[i].find(class_="date").get_text().strip()
    check_date = pd.to_datetime(date)
    print(check_date)
    #print("Colloquium",date, "\n",title,"\n",speaker,"\n",abstract,"\n",link,"\n\n")
    print("Colloquium \n",title,"\n",specs,"\n",abstract,"\n",link,"\n\n")
    title_specs = "Upcoming Event - Colloquium" + "\n" + "\n" + title + "\n" + specs + '\n'

    if ((check_date > datetime.now()) & (check_date - datetime.now() < timedelta(days = 21))):
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        left = top = width = height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = ""
    
        p = tf.add_paragraph()
        p.word_wrap = True
        p.text = title_specs
        p.font.size = Pt(26)

outreach_events = pd.read_csv('./outreach.csv')

for i in range(0,len(outreach_events['Title'])):
    #print(outreach_events['Title'][i])
    title = outreach_events['Title'][i]
    speaker = outreach_events['Speaker'][i]
    date = outreach_events['Date'][i]
    time = outreach_events['Time'][i]
    location = outreach_events['Location'][i]

    #print(date)
    check_date = pd.to_datetime(date)
    outreach_event_text = "Upcoming Event - Public Talk and Stargazing" + "\n" + "\n" + str(title) + "\n" + str(speaker) + '\n' + str(date) + ' ' + str(time) + "\n" + str(location)
    print(outreach_event_text)
    if ((check_date > datetime.now()) & (check_date - datetime.now() < timedelta(days = 28))):
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        left = top = width = height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = ""
    
        p = tf.add_paragraph()
        p.word_wrap = True
        p.text = outreach_event_text
        p.font.size = Pt(26)

now = datetime.now() # current date and time
year = now.strftime("%Y")
month = now.strftime("%m")
day = now.strftime("%d")
time = now.strftime("%H:%M:%S")
date_time = now.strftime("%Y_%m_%d_%H_%M_%S")
    
prs.save('AST_' + date_time + '.pptx')


#now for 10th floor, keeping separate for now

page = requests.get('https://www.astro.columbia.edu/events/colloquia')
contents = page.content
soup = BeautifulSoup(page.content, 'html.parser')
event = soup.find_all("a", href=lambda href: href and "event?" in href)
i=1
link = "http://www.astro.columbia.edu" + event[i].get('href')
#print(link)
page = requests.get(link)
contents = page.content
event_soup = BeautifulSoup(page.content, 'html.parser')
p = event_soup.find_all("p")
abstract = p[0].get_text()
specs = event_soup.find("ul",class_="specs").get_text().strip()
speaker = event[i].find(class_="speaker").get_text().strip()
title = event[i].find(class_="title").get_text().strip()
date = event[i].find(class_="date").get_text().strip()
#print("Colloquium \n",specs,"\n",abstract,"\n",link,"\n\n")
#print(event_soup.prettify())
title = event_soup.find_all("div",class_="content event")
title = title[-1].find_all("h3")[0].get_text()
print(title)
event = soup.find_all("a", href=lambda href: href and "event?" in href)
f = open('./calkiosk.pptx',"rb")
prs = Presentation(f)

for i in range(0,len(event)):
    link = "http://www.astro.columbia.edu" + event[i].get('href')
    #print(link)
    page = requests.get(link)
    contents = page.content
    event_soup = BeautifulSoup(page.content, 'html.parser')
    p = event_soup.find_all("p")
    abstract = p[0].get_text()
    #print("abstract:",abstract)
    #print(event_soup.prettify)
    #long_title = event[i].get('title')
    specs = event_soup.find("ul",class_="specs").get_text().strip()
    specs = specs.replace('Date:','')
    specs = specs.replace('Location:','Location: ')
    specs = specs.replace('Speaker:','')
    specs = specs.replace('Time:','')
    specs = specs.replace('Host:','Host: ')
    speaker = event[i].find(class_="speaker").get_text().strip()
    speaker = speaker.replace('Speaker:','')
    #title = event[i].find(class_="title").get_text().strip()
    title = event_soup.find_all("div",class_="content event")
    title = title[-1].find_all("h3")[0].get_text()
    date = event[i].find(class_="date").get_text().strip()
    check_date = pd.to_datetime(date)
    print(check_date)
    #print("Colloquium",date, "\n",title,"\n",speaker,"\n",abstract,"\n",link,"\n\n")
    print("Colloquium \n",title,"\n",specs,"\n",abstract,"\n",link,"\n\n")
    title_specs = "Upcoming Event - Colloquium" + "\n" + "\n" + title + "\n" + specs + '\n'

    if ((check_date > datetime.now()) & (check_date - datetime.now() < timedelta(days = 21))):
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        left = top = width = height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = ""

        p = tf.add_paragraph()
        p.word_wrap = True
        p.text = title_specs
        p.font.size = Pt(26)

for i in range(0,len(outreach_events['Title'])):
    #print(outreach_events['Title'][i])
    title = outreach_events['Title'][i]
    speaker = outreach_events['Speaker'][i]
    date = outreach_events['Date'][i]
    time = outreach_events['Time'][i]
    location = outreach_events['Location'][i]

    #print(date)
    check_date = pd.to_datetime(date)
    outreach_event_text = "Upcoming Event - Public Talk and Stargazing" + "\n" + "\n" + str(title) + "\n" + str(speaker) + '\n' + str(date) + ' ' + str(time) + "\n" + str(location)
    print(outreach_event_text)
    if ((check_date > datetime.now()) & (check_date - datetime.now() < timedelta(days = 28))):
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        left = top = width = height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = ""
    
        p = tf.add_paragraph()
        p.word_wrap = True
        p.text = outreach_event_text
        p.font.size = Pt(26)

now = datetime.now() # current date and time
year = now.strftime("%Y")
month = now.strftime("%m")
day = now.strftime("%d")
time = now.strftime("%H:%M:%S")
date_time = now.strftime("%Y_%m_%d_%H_%M_%S")
    
prs.save('CAL_' + date_time + '.pptx')

